from __future__ import annotations

import csv
import shutil
import subprocess
import tempfile
from pathlib import Path

import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


class UnsupportedSummaryError(Exception):
    def __init__(self, message: str, *, code: str = "unsupported_type") -> None:
        super().__init__(message)
        self.code = code


def _limit_text(text: str, max_chars: int = 3000) -> str:
    cleaned = text.strip()
    return cleaned[:max_chars]


def _command_path(name: str) -> str | None:
    return shutil.which(name)


def _ocr_pdf_text(file_path: Path, max_pages: int = 5) -> str:
    tesseract = _command_path("tesseract")
    if not tesseract:
        raise UnsupportedSummaryError(
            "这是图片型 PDF（扫描件），当前设备未安装 OCR 组件，暂时无法提取正文。",
            code="needs_ocr",
        )

    chunks: list[str] = []
    with tempfile.TemporaryDirectory(prefix="file-organizer-ocr-") as temp_dir:
        temp_root = Path(temp_dir)
        with fitz.open(file_path) as doc:
            page_total = min(len(doc), max_pages)
            for page_index in range(page_total):
                page = doc[page_index]
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image_path = temp_root / f"page-{page_index + 1}.png"
                pixmap.save(str(image_path))
                text = ""
                for language in ("chi_sim+eng", "eng"):
                    completed = subprocess.run(
                        [tesseract, str(image_path), "stdout", "-l", language],
                        capture_output=True,
                        text=True,
                        encoding="utf-8",
                        errors="ignore",
                        check=False,
                    )
                    text = completed.stdout.strip()
                    if text:
                        break
                if text:
                    chunks.append(text)
    combined = "\n".join(chunks).strip()
    if combined:
        return _limit_text(combined)
    raise UnsupportedSummaryError(
        "这是图片型 PDF（扫描件），已尝试 OCR，但仍未提取到可用文字。",
        code="no_text",
    )


def extract_pdf_text(file_path: Path) -> str:
    chunks: list[str] = []
    with fitz.open(file_path) as doc:
        for page_index in range(min(len(doc), 10)):
            chunks.append(doc[page_index].get_text("text"))
    combined = "\n".join(chunks).strip()
    if combined:
        return _limit_text(combined)
    return _ocr_pdf_text(file_path)


def extract_docx_text(file_path: Path) -> str:
    doc = Document(file_path)
    parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    return _limit_text("\n".join(parts), 3000)


def extract_xlsx_text(file_path: Path) -> str:
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"Sheet: {sheet.title}")
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > 20:
                break
            values = [str(cell) if cell is not None else "" for cell in row]
            parts.append(" | ".join(values))
    workbook.close()
    return "\n".join(parts).strip()


def extract_pptx_text(file_path: Path) -> str:
    presentation = Presentation(file_path)
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and str(shape.text).strip():
                parts.append(str(shape.text).strip())
    return "\n".join(parts).strip()


def _extract_legacy_office_text(file_path: Path) -> str:
    soffice = _command_path("soffice")
    if not soffice:
        raise UnsupportedSummaryError(
            f"这是旧版 Office 文件（{file_path.suffix.lower()}），当前设备未安装 LibreOffice，暂时无法自动转换。请先另存为新版格式后再整理。",
            code="needs_conversion",
        )

    suffix_map = {
        ".doc": "docx",
        ".xls": "xlsx",
        ".ppt": "pptx",
    }
    target_ext = suffix_map.get(file_path.suffix.lower())
    if not target_ext:
        raise UnsupportedSummaryError(
            f"暂不支持此文件类型的摘要提取：{file_path.suffix.lower()}",
            code="unsupported_type",
        )

    with tempfile.TemporaryDirectory(prefix="file-organizer-office-") as temp_dir:
        output_dir = Path(temp_dir)
        completed = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                target_ext,
                "--outdir",
                str(output_dir),
                str(file_path),
            ],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="ignore",
            check=False,
        )
        converted = output_dir / f"{file_path.stem}.{target_ext}"
        if not converted.exists():
            detail = completed.stderr.strip() or completed.stdout.strip() or "转换工具未返回详细信息。"
            raise UnsupportedSummaryError(
                f"旧版 Office 文件自动转换失败：{detail}",
                code="needs_conversion",
            )
        if target_ext == "docx":
            return extract_docx_text(converted)
        if target_ext == "xlsx":
            return extract_xlsx_text(converted)
        if target_ext == "pptx":
            return extract_pptx_text(converted)
    raise UnsupportedSummaryError(
        f"暂不支持此文件类型的摘要提取：{file_path.suffix.lower()}",
        code="unsupported_type",
    )


def extract_text_text(file_path: Path) -> str:
    return _limit_text(file_path.read_text(encoding="utf-8", errors="ignore"), 3000)


def extract_csv_text(file_path: Path) -> str:
    max_rows = 200
    max_chars = 3000
    collected: list[str] = []
    total_chars = 0
    with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for index, row in enumerate(reader, start=1):
            if index > max_rows:
                break
            line = ", ".join(row).strip()
            if not line:
                continue
            remaining = max_chars - total_chars
            if remaining <= 0:
                break
            if len(line) > remaining:
                line = line[:remaining]
            collected.append(line)
            total_chars += len(line) + 1
            if total_chars >= max_chars:
                break
    return "\n".join(collected).strip()


def extract_text(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".docx":
        return extract_docx_text(path)
    if suffix == ".xlsx":
        return extract_xlsx_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    if suffix in {".txt", ".md"}:
        return extract_text_text(path)
    if suffix == ".csv":
        return extract_csv_text(path)
    if suffix in {".doc", ".xls", ".ppt"}:
        return _extract_legacy_office_text(path)
    raise UnsupportedSummaryError(
        f"暂不支持此文件类型的摘要提取：{suffix}",
        code="unsupported_type",
    )
